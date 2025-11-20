#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
提取PPT内容，并为每张图片生成有意义的文件名
"""

import json
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from zipfile import ZipFile
import shutil
import re

def extract_text_from_shape(shape):
    """从形状中提取文本"""
    if not shape.has_text_frame:
        return None
    text_parts = []
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip():
                text_parts.append(run.text.strip())
    return '\n'.join(text_parts) if text_parts else None

def clean_filename(text):
    """清理文本，生成有效的文件名"""
    if not text:
        return "image"
    
    # 提取关键词
    text = text.replace('\n', ' ')
    text = re.sub(r'[^\w\s-]', '', text)
    text = re.sub(r'\s+', '_', text)
    
    # 限制长度
    if len(text) > 50:
        words = text.split('_')
        text = '_'.join(words[:5])
    
    return text.lower() or "image"

def extract_slide_content_with_images(slide, slide_num, prs):
    """提取幻灯片内容，包括图片信息"""
    content = {
        'slide_number': slide_num,
        'texts': [],
        'images': [],
        'title': ''
    }
    
    shape_count = 0
    for shape in slide.shapes:
        # 提取文本
        if shape.has_text_frame:
            text = extract_text_from_shape(shape)
            if text:
                content['texts'].append(text)
                # 如果文本较短且看起来像标题，记录为标题
                if len(text) < 50 and len(content['title']) == 0:
                    content['title'] = text
    
    return content

def organize_images_by_slide(prs, lang='en'):
    """按幻灯片组织图片"""
    image_map = []
    
    for i, slide in enumerate(prs.slides, 1):
        slide_images = []
        shape_idx = 0
        
        # 先提取文本作为上下文
        context_texts = []
        for para_shape in slide.shapes:
            if para_shape.has_text_frame:
                text = extract_text_from_shape(para_shape)
                if text and len(text) < 50:
                    context_texts.append(text)
        
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shape_idx += 1
                
                # 使用第一个短文本作为上下文
                context_text = context_texts[0] if context_texts else ""
                
                slide_images.append({
                    'shape_index': shape_idx,
                    'context': context_text,
                    'slide_num': i
                })
        
        image_map.append({
            'slide': i,
            'images': slide_images,
            'title': extract_slide_content_with_images(slide, i, prs)['title']
        })
    
    return image_map

def extract_and_rename_images(pptx_path, lang='en', output_dir=None):
    """提取并重命名图片"""
    from pptx import Presentation
    
    print(f"\n正在提取图片并生成有意义的文件名...")
    
    # 读取PPT获取图片上下文
    prs = Presentation(str(pptx_path))
    image_map = organize_images_by_slide(prs, lang)
    
    # 提取图片
    if not output_dir:
        output_dir = Path(f'extracted_content/{lang}_images_renamed')
    output_dir.mkdir(parents=True, exist_ok=True)
    
    extracted_images = []
    
    with ZipFile(pptx_path, 'r') as zip_ref:
        all_files = zip_ref.namelist()
        media_files = [f for f in all_files if f.startswith('ppt/media/')]
        
        # 按数字排序
        media_files_sorted = sorted(media_files, key=lambda x: int(re.findall(r'\d+', Path(x).stem)[0]) if re.findall(r'\d+', Path(x).stem) else 999999)
        
        image_idx = 0
        for slide_info in image_map:
            for img_info in slide_info['images']:
                if image_idx < len(media_files_sorted):
                    old_file = media_files_sorted[image_idx]
                    
                    # 获取文件扩展名
                    ext = Path(old_file).suffix
                    
                    # 生成新文件名
                    if slide_info['title']:
                        new_name = clean_filename(slide_info['title'])
                        new_name = f"slide{slide_info['slide']:02d}_{new_name}_{img_info['shape_index']}{ext}"
                    elif img_info['context']:
                        new_name = f"slide{slide_info['slide']:02d}_{clean_filename(img_info['context'])}_{img_info['shape_index']}{ext}"
                    else:
                        new_name = f"slide{slide_info['slide']:02d}_image{img_info['shape_index']}{ext}"
                    
                    # 复制并重命名
                    output_path = output_dir / new_name
                    with zip_ref.open(old_file) as source, open(output_path, 'wb') as target:
                        shutil.copyfileobj(source, target)
                    
                    extracted_images.append({
                        'old_name': Path(old_file).name,
                        'new_name': new_name,
                        'slide': slide_info['slide'],
                        'context': slide_info['title']
                    })
                    
                    image_idx += 1
    
    # 保存映射关系
    mapping_file = output_dir / 'image_mapping.json'
    with open(mapping_file, 'w', encoding='utf-8') as f:
        json.dump(extracted_images, f, ensure_ascii=False, indent=2)
    
    print(f"✓ 提取了 {len(extracted_images)} 张图片")
    print(f"✓ 已保存映射文件: {mapping_file}")
    
    return extracted_images

def process_ppt(pptx_path, lang='en'):
    """处理单个PPT文件，包含图片重命名"""
    pptx_path = Path(pptx_path)
    print(f"\n{'='*60}")
    print(f"处理: {pptx_path.name} ({lang})")
    print(f"{'='*60}")
    
    if not pptx_path.exists():
        print(f"❌ 文件不存在: {pptx_path}")
        return None
    
    # 创建输出目录
    output_base = Path('extracted_content')
    output_base.mkdir(exist_ok=True)
    
    # 提取并重命名图片
    print(f"📸 正在提取并重命名图片...")
    extracted_images = extract_and_rename_images(pptx_path, lang)
    images_dir = output_base / f'{lang}_images_renamed'
    
    # 读取PPT内容
    print(f"\n📄 正在读取PPT内容...")
    prs = Presentation(str(pptx_path))
    
    slides_data = []
    for i, slide in enumerate(prs.slides, 1):
        print(f"  处理幻灯片 {i}/{len(prs.slides)}...", end='\r')
        
        slide_content = extract_slide_content_with_images(slide, i, prs)
        
        # 关联重命名后的图片
        slide_content['images_renamed'] = [
            img for img in extracted_images 
            if img['slide'] == i
        ]
        
        slides_data.append(slide_content)
    
    print(f"\n✓ 处理了 {len(slides_data)} 张幻灯片")
    
    # 保存数据
    data_file = output_base / f'{lang}_slides_with_named_images.json'
    with open(data_file, 'w', encoding='utf-8') as f:
        json.dump(slides_data, f, ensure_ascii=False, indent=2)
    
    print(f"✓ 数据已保存到: {data_file}")
    
    return {
        'lang': lang,
        'slides': slides_data,
        'images': extracted_images,
        'total_slides': len(slides_data),
        'images_dir': images_dir,
        'mapping_file': images_dir / 'image_mapping.json'
    }

def organize_by_pages_with_named_images(en_data, cn_data):
    """按页面组织内容，使用重命名后的图片"""
    print(f"\n{'='*60}")
    print("正在组织页面结构（使用有意义的文件名）...")
    print(f"{'='*60}")
    
    pages_dir = Path('website_data_named')
    pages_dir.mkdir(exist_ok=True)
    
    max_slides = max(en_data['total_slides'], cn_data['total_slides'])
    
    for i in range(max_slides):
        page_dir = pages_dir / f'page_{i+1:02d}'
        page_dir.mkdir(exist_ok=True)
        
        # 创建子目录
        page_dir.joinpath('texts').mkdir(exist_ok=True)
        page_dir.joinpath('images').mkdir(exist_ok=True)
        
        # 处理英文内容
        if i < en_data['total_slides']:
            en_slide = en_data['slides'][i]
            
            # 保存文本
            texts_file = page_dir / 'texts' / 'en.txt'
            with open(texts_file, 'w', encoding='utf-8') as f:
                for idx, text in enumerate(en_slide['texts']):
                    f.write(f"--- Text Block {idx+1} ---\n")
                    f.write(text + "\n\n")
            
            # 复制重命名后的图片
            if 'images_renamed' in en_slide:
                for img_info in en_slide['images_renamed']:
                    source = en_data['images_dir'] / img_info['new_name']
                    if source.exists():
                        dest = page_dir / 'images' / img_info['new_name']
                        shutil.copy2(source, dest)
        
        # 处理中文内容
        if i < cn_data['total_slides']:
            cn_slide = cn_data['slides'][i]
            
            # 保存文本
            texts_file = page_dir / 'texts' / 'cn.txt'
            with open(texts_file, 'w', encoding='utf-8') as f:
                for idx, text in enumerate(cn_slide['texts']):
                    f.write(f"--- 文本块 {idx+1} ---\n")
                    f.write(text + "\n\n")
            
            # 复制重命名后的图片
            if 'images_renamed' in cn_slide:
                for img_info in cn_slide['images_renamed']:
                    source = cn_data['images_dir'] / img_info['new_name']
                    if source.exists():
                        dest = page_dir / 'images' / img_info['new_name']
                        shutil.copy2(source, dest)
        
        print(f"  ✓ {page_dir}")
    
    return pages_dir

def main():
    print("🚀 开始提取PPT内容（图片将使用有意义的文件名）...")
    print("="*60)
    
    # 处理英文PPT
    en_data = process_ppt('3amClub EN.pptx', 'en')
    
    # 处理中文PPT
    cn_data = process_ppt('3amClub2024.pptx', 'cn')
    
    if not en_data or not cn_data:
        print("❌ 处理失败")
        return
    
    # 组织页面结构
    pages_dir = organize_by_pages_with_named_images(en_data, cn_data)
    
    print("\n" + "="*60)
    print("✅ 完成！")
    print("="*60)
    print(f"📁 页面数据: {pages_dir.absolute()}")
    print(f"📊 总共 {len(list(pages_dir.glob('page_*')))} 个页面文件夹")
    print(f"🖼️  图片已重命名为有意义的名称")
    print(f"\n📝 图片映射信息:")
    print(f"  • 英文: {en_data['mapping_file']}")
    print(f"  • 中文: {cn_data['mapping_file']}")

if __name__ == "__main__":
    main()


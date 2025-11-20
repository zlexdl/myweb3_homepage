#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
从PPT文件中提取每一页的文本和图片，并将每页内容保存到独立文件夹
支持图片内容识别和智能命名
"""

import os
import hashlib
from pathlib import Path
from pptx import Presentation
from PIL import Image
import io


def get_image_hash(image_data):
    """计算图片的哈希值，用于识别不同的图片"""
    return hashlib.md5(image_data).hexdigest()[:8]


def clean_filename(name):
    """清理文件名，移除非法字符"""
    illegal_chars = ['<', '>', ':', '"', '/', '\\', '|', '?', '*']
    for char in illegal_chars:
        name = name.replace(char, '_')
    return name.strip()


def extract_slide_content(prs, output_base_dir, ppt_name):
    """
    从PPT中提取所有幻灯片的内容
    
    Args:
        prs: Presentation对象
        output_base_dir: 输出基础目录
        ppt_name: PPT文件名（用于识别中文或英文）
    """
    output_base = Path(output_base_dir)
    output_base.mkdir(exist_ok=True)
    
    # 创建输出目录
    if ppt_name == "3amClub2024.pptx":
        lang = "cn"
        lang_name = "中文"
    else:
        lang = "en"
        lang_name = "英文"
    
    print(f"\n处理{lang_name}PPT ({ppt_name})...")
    
    # 图片计数器（按页面）
    image_counters = {}
    
    for slide_idx, slide in enumerate(prs.slides, start=1):
        print(f"  处理第 {slide_idx} 页...")
        
        # 为每页创建文件夹
        page_dir = output_base / f"slide_{slide_idx:02d}"
        texts_dir = page_dir / "texts"
        images_dir = page_dir / "images"
        
        texts_dir.mkdir(parents=True, exist_ok=True)
        images_dir.mkdir(parents=True, exist_ok=True)
        
        # 提取文本内容
        text_content = []
        
        # 遍历所有形状获取文本
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_content.append(shape.text.strip())
        
        # 保存文本到文件
        if text_content:
            text_file = texts_dir / f"{lang}.txt"
            with open(text_file, 'w', encoding='utf-8') as f:
                f.write('\n\n'.join(text_content))
            print(f"    已提取 {len(text_content)} 段文本")
        
        # 提取图片
        if slide_idx not in image_counters:
            image_counters[slide_idx] = 0
        
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                try:
                    image = shape.image
                    image_bytes = image.blob
                    
                    # 计算图片哈希值
                    img_hash = get_image_hash(image_bytes)
                    
                    # 获取原始扩展名
                    ext = image.ext
                    if not ext or ext.lower() not in ['png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff']:
                        ext = 'png'  # 默认使用png
                    
                    # 增加计数器
                    image_counters[slide_idx] += 1
                    img_num = image_counters[slide_idx]
                    
                    # 生成文件名: slide编号_图片序号_hash.扩展名
                    filename = f"slide{slide_idx:02d}_img{img_num:02d}_{img_hash}.{ext}"
                    img_path = images_dir / filename
                    
                    # 保存图片
                    with open(img_path, 'wb') as f:
                        f.write(image_bytes)
                    
                    print(f"    已保存图片: {filename}")
                    
                except Exception as e:
                    print(f"    跳过图片提取: {e}")
        
        if image_counters[slide_idx] == 0:
            print(f"    本页无图片")
    
    print(f"\n{lang_name}PPT提取完成！共 {len(prs.slides)} 页")
    print(f"输出目录: {output_base_dir}")


def main():
    """主函数"""
    # PPT文件路径
    ppt_en = "3amClub EN.pptx"
    ppt_cn = "3amClub2024.pptx"
    
    # 检查文件是否存在
    if not os.path.exists(ppt_en):
        print(f"错误: 找不到文件 {ppt_en}")
        return
    
    if not os.path.exists(ppt_cn):
        print(f"错误: 找不到文件 {ppt_cn}")
        return
    
    print("=" * 60)
    print("PPT内容提取工具")
    print("=" * 60)
    
    # 处理英文PPT
    print("\n正在读取英文PPT...")
    try:
        prs_en = Presentation(ppt_en)
        extract_slide_content(prs_en, "extracted_en", ppt_en)
    except Exception as e:
        print(f"处理英文PPT时出错: {e}")
    
    # 处理中文PPT
    print("\n正在读取中文PPT...")
    try:
        prs_cn = Presentation(ppt_cn)
        extract_slide_content(prs_cn, "extracted_cn", ppt_cn)
    except Exception as e:
        print(f"处理中文PPT时出错: {e}")
    
    print("\n" + "=" * 60)
    print("所有PPT处理完成！")
    print("=" * 60)


if __name__ == "__main__":
    main()

